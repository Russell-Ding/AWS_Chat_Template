from flask import Flask, render_template, request, jsonify
import database as db
import boto3
import json
import os
import requests
from bs4 import BeautifulSoup
from werkzeug.utils import secure_filename
import base64
import mimetypes
import pypdf
import docx
import openpyxl
import pptx

app = Flask(__name__)

# --- File Upload and Content Configuration ---
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'txt', 'json', 'pdf', 'doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx', 'png', 'jpg', 'jpeg', 'gif', 'webp'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def read_file_content(filepath):
    """Reads content from various file types, encoding images to base64."""
    try:
        mime_type, _ = mimetypes.guess_type(filepath)
        if mime_type and mime_type.startswith('image/'):
            with open(filepath, "rb") as f:
                return base64.b64encode(f.read()).decode('utf-8'), mime_type
        
        text_content = ""
        ext = filepath.rsplit('.', 1)[1].lower()

        if ext == 'pdf':
            with open(filepath, 'rb') as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    text_content += page.extract_text() or ''
        elif ext == 'docx':
            doc = docx.Document(filepath)
            for para in doc.paragraphs:
                text_content += para.text + '\n'
        elif ext == 'xlsx':
            workbook = openpyxl.load_workbook(filepath)
            for sheet in workbook.sheetnames:
                worksheet = workbook[sheet]
                for row in worksheet.iter_rows(values_only=True):
                    text_content += '\t'.join([str(cell) if cell is not None else '' for cell in row]) + '\n'
        elif ext == 'pptx':
            prs = pptx.Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + '\n'
        else: # Default to plain text
            with open(filepath, 'r', encoding='utf-8') as f:
                text_content = f.read()
        
        return text_content, mime_type or 'text/plain'

    except Exception as e:
        return f"Error reading file {os.path.basename(filepath)}: {e}", "text/plain"

# --- Tool Definition: Google Search and Multi-Scrape ---
def google_search(query):
    # ... [The existing google_search function remains unchanged] ...
    pass # Placeholder for brevity

@app.route("/")
def index():
    conversations = db.get_conversations()
    return render_template("index.html", conversations=conversations)

@app.route("/upload", methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    file = request.files['file']
    if file.filename == '' or not allowed_file(file.filename):
        return jsonify({"error": "File type not allowed"}), 400
    if file:
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        return jsonify({"filepath": filepath}), 200

@app.route("/conversation/<int:conversation_id>")
def get_conversation_route(conversation_id):
    conversation = db.get_conversation(conversation_id)
    if conversation:
        return jsonify(conversation)
    return jsonify({"error": "Conversation not found"}), 404

@app.route("/chat", methods=["POST"])
def chat():
    data = request.json
    message_text = data.get("message")
    model = data.get("model")
    conversation_id = data.get("conversation_id")
    files = data.get("files", [])
    new_conversation_info = None

    if not conversation_id:
        name = message_text[:20] if message_text else "New Conversation"
        conversation_id = db.create_conversation(name=f"Conversation about {name}...", model=model)
        new_conversation_info = {"id": conversation_id, "name": name}

    # --- Construct User Message with File Content ---
    user_message_content = []
    if message_text:
        user_message_content.append({"type": "text", "text": message_text})

    for fpath in files:
        content, mime_type = read_file_content(fpath)
        if mime_type.startswith('image/'):
            user_message_content.append({
                "type": "image",
                "source": {"type": "base64", "media_type": mime_type, "data": content}
            })
        else:
            # Prepend text content to the user's message for context
            text_for_prompt = f"\n\n--- Content from {os.path.basename(fpath)} ---\n{content}"
            # Add this as a separate message or prepend to the main one
            if not message_text and not user_message_content:
                 user_message_content.append({"type": "text", "text": text_for_prompt})
            else:
                user_message_content[0]["text"] += text_for_prompt

    db.add_message(conversation_id, "user", json.dumps(user_message_content)) # Store structured content

    llm_response = ""
    try:
        bedrock_runtime = boto3.client(service_name='bedrock-runtime')
        max_turns = 3
        for turn in range(max_turns):
            conversation_history = db.get_conversation(conversation_id)["messages"]
            
            # Reformat history for model
            formatted_history = []
            for msg in conversation_history:
                try:
                    content_data = json.loads(msg['content'])
                except json.JSONDecodeError:
                    content_data = [{"type": "text", "text": msg['content']}]
                formatted_history.append({"role": msg['role'], "content": content_data})

            system_prompt = '''You are a helpful assistant. If you need to find out recent information or anything you don't know, you can use the google_search tool. To use it, you MUST respond with ONLY a JSON object containing 'tool_name': 'google_search' and 'query': 'your search query'. Do not add any other text or explanation.'''

            body = ""
            if "anthropic" in model:
                body = json.dumps({
                    "anthropic_version": "bedrock-2023-05-31",
                    "max_tokens": 4096,
                    "system": system_prompt,
                    "messages": formatted_history
                })
            else:
                # Simplified logic for non-Anthropic models
                prompt = "\n\n".join([f["{m['role']}"]: {m['content'][0]['text']}" for m in formatted_history])
                body = json.dumps({"prompt": prompt, "max_tokens": 4096})

            response = bedrock_runtime.invoke_model(
                body=body, modelId=model, accept="application/json", contentType="application/json"
            )
            response_body = json.loads(response.get("body").read())

            if "anthropic" in model:
                llm_response = response_body.get('content', [{}])[0].get('text', "")
            else:
                llm_response = response_body.get('completion', str(response_body))

            # --- Robust Tool Use Check (BUG FIXED) ---
            try:
                start_index = llm_response.find('{')
                end_index = llm_response.rfind('}') + 1
                if start_index != -1 and end_index != -1:
                    json_str = llm_response[start_index:end_index]
                    tool_call = json.loads(json_str)
                    if tool_call.get("tool_name") == "google_search":
                        search_query = tool_call.get("query")
                        search_results = google_search(search_query)
                        db.add_message(conversation_id, "assistant", llm_response)
                        db.add_message(conversation_id, "user", f"Search results for \"{search_query}\": {search_results}")
                        continue
                break
            except (json.JSONDecodeError, AttributeError):
                break

    except Exception as e:
        llm_response = f"Error communicating with Bedrock: {e}"

    db.add_message(conversation_id, "assistant", llm_response)

    response_data = {
        "conversation_id": conversation_id,
        "messages": db.get_conversation(conversation_id)["messages"],
    }
    if new_conversation_info:
        response_data["new_conversation"] = new_conversation_info

    return jsonify(response_data)

if __name__ == "__main__":
    app.run(debug=True)
